#!/usr/bin/env python3
r"""
D:\Krafton\ 문서를 마크다운 텍스트로 추출하는 스크립트.
PPTX -> python-pptx, PDF -> PyPDF2 사용.
출력: Docs/V2/Research/Ingested/<category>/<slug>.md
"""

import argparse
import os
import sys
import re
from pathlib import Path
from datetime import datetime

# ── 설정 ──────────────────────────────────────────────
PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_ROOT = PROJECT_ROOT / "Docs" / "V2" / "Research" / "Ingested"

# 추출 대상 문서 매핑: (원본 경로 패턴, 카테고리, 출력 slug)
DOCUMENT_MAP = [
    # PD 방법론
    ("PD Guild Books - KRAFTON Producer Guild - Wiki.pdf", "methodology", "pd-guild-books"),
    ("KRAFTON PD직군_게임산업OnBoarding_박민현_2106_r.pptx", "methodology", "krafton-pd-onboarding"),
    ("Krafton_PD_Pathfinders 2023 OnBoarding.pptx", "methodology", "krafton-pathfinders-2023-onboarding"),
    ("PDGuild_TeamCulture_2106.pptx", "methodology", "pdguild-team-culture"),
    ("게임의 조건_박현규.pptx", "methodology", "game-conditions-parkhyeongyu"),
    ("재미와 디자인과 개발에 대한 저의 생각 - 이종인 (Jongin Lee) - Wiki.pdf", "methodology", "fun-design-dev-jonginlee"),
    ("제작 방향성 - Memento Mori - Wiki-압축됨.pdf", "methodology", "memento-mori-production-direction"),
    ("더 포텐셜 제안 플로우 - Memento Mori - Wiki.pdf", "methodology", "memento-mori-potential-flow"),
    ("Krafton_PD_Pathfinders 2기 채용설명회_V0.8.pptx", "methodology", "krafton-pathfinders-2-recruitment"),
    ("PDGuild KraftonPDGuild 공채 OnBoarding_V0.2.pptx", "methodology", "pdguild-onboarding-v02"),
    # 시장조사
    ("2nd Project 조사 방법론 - Publishing Dept. 5 - Wiki.pdf", "market-research", "2nd-project-methodology"),
    ("2nd 프로젝트 조사 자료 - Publishing Dept. 5 - Wiki.pdf", "market-research", "2nd-project-data"),
    ("Salesforce - Trend Slack 2024 - Publishing Dept. 5 - Wiki.pdf", "market-research", "salesforce-trend-2024"),
    ("240229_2023 쿨재팬마켓 오타쿠시장의 철저 연구_YanoResearch_(Kor 번역).pdf", "market-research", "cool-japan-otaku-market-2023"),
    ("돈을 쓰고 돈을 버는 퍼블리싱의 본질 (The essence of the game publishing business, where money is spent and earned) - Publishing Dept. 5 - Wiki.pdf", "market-research", "publishing-essence"),
    # 포스트모템
    ("Project D Postmortem - Game Production_공유 - Wiki.pdf", "postmortem", "project-d-postmortem"),
    ("2024-01-26 정글 게임 랩 1기 포스트모템 by 임성진.pptx", "postmortem", "jungle-gamelab-postmortem"),
    ("Postmortem_임성진 (1).pptx", "postmortem", "postmortem-lim-seongjin"),
    ("[Pathfinders Team] 신입 공채 Product Research - CaseStudy 수행 과정 Note - KRAFTON Producer Guild - Wiki.pdf", "postmortem", "pathfinders-case-study"),
    # 포스트모템 - BRO 시리즈
    ("CH/BRO - 2017년 하반기.pdf", "postmortem", "bro-2017-h2"),
    ("CH/BRO Sprint 1.pdf", "postmortem", "bro-sprint-1"),
    ("CH/BRO_2016_Workshop.pdf", "postmortem", "bro-2016-workshop"),
    ("CH/Milestone 2 Review_160811 (revised).pdf", "postmortem", "bro-milestone-2-review"),
    ("CH/Milestone 3 Review.pdf", "postmortem", "bro-milestone-3-review"),
    ("CH/Milestone 4 Review_Team.pdf", "postmortem", "bro-milestone-4-review"),
    ("CH/Project BRO (team).pdf", "postmortem", "project-bro-team"),
    ("CH/TSL-pdf-mid.pdf", "postmortem", "tsl-mid"),
    # 레퍼런스
    ("PMO/Chan/Brawlstars_Mod_Count.xlsx", "reference", "brawlstars-mod-count"),
    # 대형 PPTX
    ("PMO/IdeationTF_김우찬_TSS_Final.pptx", "methodology", "ideationtf-tss-final"),
    ("kizna/ProjectC3_Proposal_final_2차발표.pptx", "reference", "projectc3-proposal"),
    ("kizna/ProjectC3_아트제작방향성+AD후보리스트.pptx", "reference", "projectc3-art-direction"),
    # 기타 방법론
    ("2nd Project 개발역량강화 A안 v1.0 - Publishing Dept. 5 - Wiki.pdf", "methodology", "2nd-project-dev-capability"),
    ("Project-based org. - Publishing Dept. 5 - Wiki.pdf", "methodology", "project-based-org"),
    ("[Match Making] MMR 시스템 기획서 (v 2025.1) - PUBG PC&Console - Wiki.pdf", "reference", "pubg-mmr-system"),
    ("[국문] THE CREATIVE_ 제안 가이드라인 - PUBG IP FRANCHISE IDEATION TF - Wiki.pdf", "reference", "pubg-creative-guideline"),
    # ── Ref_Dinkum (딩컴 모바일 기획 자료) ──
    ("Ref_Dinkum/01. 딩컴 모바일 제작 방향_1 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-direction-1"),
    ("Ref_Dinkum/01. 딩컴 모바일 제작 방향_2 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-direction-2"),
    ("Ref_Dinkum/01. 딩컴 모바일 제작 방향_3 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-direction-3"),
    ("Ref_Dinkum/01. 딩컴 모바일 제작 방향_4 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-direction-4"),
    ("Ref_Dinkum/01. 딩컴 모바일 제작 방향_5 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-direction-5"),
    ("Ref_Dinkum/01. 딩컴 모바일 제작 방향_6 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-direction-6"),
    ("Ref_Dinkum/01. 딩컴 모바일 제작 방향_7 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-direction-7"),
    ("Ref_Dinkum/01. 딩컴 모바일 제작 방향_8 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-direction-8"),
    ("Ref_Dinkum/01. 딩컴 모바일 제작 방향_9 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-direction-9"),
    ("Ref_Dinkum/AI 구조 문서 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-ai-architecture"),
    ("Ref_Dinkum/AI 데이터 설명 문서 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-ai-data"),
    ("Ref_Dinkum/[3차 발주] G-star 전후 발주 아이템 정리 페이지 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-gstar-items"),
    ("Ref_Dinkum/[PC] 의상 장착(룩북) 시스템 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-lookbook-system"),
    ("Ref_Dinkum/[PC] 캐릭터 생성 시스템 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-character-creation"),
    ("Ref_Dinkum/개발용 치트 기능 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-cheat-system"),
    ("Ref_Dinkum/게임 컨셉 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-game-concept"),
    ("Ref_Dinkum/공용 UI 시스템 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-ui-system"),
    ("Ref_Dinkum/기존 전투 시스템 분석 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-combat-analysis"),
    ("Ref_Dinkum/기획서 템플릿 - 5minlab (5ml) - Wiki.pdf", "methodology", "dinkum-gdd-template"),
    ("Ref_Dinkum/무기 데이터 구조 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-weapon-data"),
    ("Ref_Dinkum/새로운 전투 시스템 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-new-combat"),
    ("Ref_Dinkum/아이템별 Hand 상태 모델링 정리 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-hand-state"),
    ("Ref_Dinkum/에셋 관리 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-asset-management"),
    ("Ref_Dinkum/유저 시나리오 초안 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-user-scenarios"),
    ("Ref_Dinkum/추가된 데이터 별 확인 절차(QA) - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-qa-procedure"),
    ("Ref_Dinkum/타일 오브젝트 배치 시스템 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-tile-placement"),
    ("Ref_Dinkum/플레이어섬 리젠 규칙 - 5minlab (5ml) - Wiki.pdf", "reference", "dinkum-island-regen"),
    # ── Ref_HotDrop (PUBG 랜드마크/핫드랍 모드) ──
    ("Ref_HotDrop/#33.1 랜드마크 모드 밸런스 폴리싱 (블루존 + 아이템 스폰) - PUBG PC&Console - Wiki.pdf", "reference", "hotdrop-balance-polishing"),
    ("Ref_HotDrop/#34.2 핫드랍 모드 스폰 롤백 - PUBG PC&Console - Wiki.pdf", "reference", "hotdrop-spawn-rollback"),
    ("Ref_HotDrop/(old) 대도시 꼬라박기 모드 기획 검토 - PUBG PC&Console - Wiki.pdf", "reference", "hotdrop-old-mode-review"),
    ("Ref_HotDrop/1. Landmark - Information - PUBG PC&Console - Wiki.pdf", "reference", "hotdrop-landmark-info"),
    ("Ref_HotDrop/10. Landmark - Postmortem - PUBG PC&Console - Wiki.pdf", "postmortem", "hotdrop-landmark-postmortem"),
    ("Ref_HotDrop/240830 플레이데이 랜드마크 모드 플레이데이 블루존 시나리오 문서 - PUBG PC&Console - Wiki.pdf", "reference", "hotdrop-playday-bluezone"),
    ("Ref_HotDrop/3. Landmark - Task - PUBG PC&Console - Wiki.pdf", "reference", "hotdrop-landmark-task"),
    ("Ref_HotDrop/Landmark - Service Report - PUBG PC&Console - Wiki.pdf", "postmortem", "hotdrop-service-report"),
    ("Ref_HotDrop/랜드마크 모드 - 블루존 밸런스 기획 - PUBG PC&Console - Wiki.pdf", "reference", "hotdrop-bluezone-balance"),
    ("Ref_HotDrop/랜드마크 모드 - 시스템 기획서 - PUBG PC&Console - Wiki.pdf", "reference", "hotdrop-system-design"),
    ("Ref_HotDrop/랜드마크 모드 - 전장 기획서 - PUBG PC&Console - Wiki.pdf", "reference", "hotdrop-battlefield-design"),
    ("Ref_HotDrop/랜드마크 모드 - 키 아트 요청서 - PUBG PC&Console - Wiki.pdf", "reference", "hotdrop-key-art-request"),
    ("Ref_HotDrop/플레이데이 블루존 & 라운드 시간 관련 피드백 요약 - PUBG PC&Console - Wiki.pdf", "reference", "hotdrop-playday-feedback"),
    # ── PMO 추가 ──
    ("PMO/(CD) 240716_Certain Affinity_Paul Sams.pdf", "reference", "pmo-certain-affinity"),
    ("PMO/2023_BRP_draft_Pixel_Share with BX_l (1).pdf", "reference", "pmo-brp-draft"),
    ("PMO/Chan/TD.pptx", "reference", "pmo-td"),
    ("PMO/Chan/TDS.pptx", "reference", "pmo-tds"),
    # ── kizna 추가 ──
    ("kizna/AllHands_KIZNA Cell_240822(KR)_Final.pdf", "reference", "kizna-allhands"),
    # ── Jungle 추가 (게임 관련만) ──
    ("Jungle/Appliance/250110_TSS.pptx", "reference", "jungle-tss"),
]

# 개인 문서 제외 패턴
EXCLUDE_PATTERNS = [
    "여권", "국세청", "고용보험", "구직등록", "월급날", "급여명세",
    "수료증", "취업", "KakaoTalk"
]

KRAFTON_ROOT = Path("D:/Krafton")


def slugify(text: str) -> str:
    """텍스트를 파일명용 slug로 변환"""
    text = re.sub(r'[^\w\s가-힣-]', '', text)
    text = re.sub(r'\s+', '-', text.strip())
    return text.lower()[:80]


def extract_pdf(filepath: Path, max_pages: int = 0) -> tuple[str, dict]:
    """PDF에서 텍스트 추출. 반환: (텍스트, 메타데이터)"""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return "", {"error": "PyPDF2 not installed"}

    try:
        reader = PdfReader(str(filepath))
        total_pages = len(reader.pages)
        pages_to_read = total_pages if max_pages <= 0 else min(max_pages, total_pages)

        text_parts = []
        for i in range(pages_to_read):
            page_text = reader.pages[i].extract_text() or ""
            if page_text.strip():
                text_parts.append(f"## Page {i + 1}\n\n{page_text.strip()}")

        meta = {
            "total_pages": total_pages,
            "extracted_pages": pages_to_read,
            "format": "PDF",
        }
        return "\n\n---\n\n".join(text_parts), meta
    except Exception as e:
        return "", {"error": str(e)}


def extract_pptx(filepath: Path, max_slides: int = 0) -> tuple[str, dict]:
    """PPTX에서 텍스트 추출. 반환: (텍스트, 메타데이터)"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return "", {"error": "python-pptx not installed"}

    try:
        prs = Presentation(str(filepath))
        total_slides = len(prs.slides)
        slides_to_read = total_slides if max_slides <= 0 else min(max_slides, total_slides)

        text_parts = []
        for i, slide in enumerate(prs.slides):
            if i >= slides_to_read:
                break

            slide_texts = []
            slide_title = ""

            for shape in slide.shapes:
                # 제목 추출
                if shape.has_text_frame:
                    if shape.shape_id == 0 or (hasattr(shape, 'placeholder_format') and
                            shape.placeholder_format is not None and
                            shape.placeholder_format.idx == 0):
                        slide_title = shape.text_frame.text.strip()

                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            slide_texts.append(para_text)

                # 표 추출
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    if rows:
                        # 마크다운 테이블 형식
                        header = rows[0]
                        separator = " | ".join(["---"] * len(rows[0].split(" | ")))
                        table_md = f"\n{header}\n{separator}\n" + "\n".join(rows[1:])
                        slide_texts.append(table_md)

            # 노트 추출
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            title_line = f"## Slide {i + 1}: {slide_title}" if slide_title else f"## Slide {i + 1}"
            content = "\n".join(slide_texts) if slide_texts else "(빈 슬라이드 또는 이미지만 포함)"
            notes = f"\n\n> **Notes**: {notes_text}" if notes_text else ""

            text_parts.append(f"{title_line}\n\n{content}{notes}")

        meta = {
            "total_slides": total_slides,
            "extracted_slides": slides_to_read,
            "format": "PPTX",
        }
        return "\n\n---\n\n".join(text_parts), meta
    except Exception as e:
        return "", {"error": str(e)}


def extract_xlsx(filepath: Path) -> tuple[str, dict]:
    """XLSX에서 텍스트 추출 (간단한 CSV 변환)"""
    try:
        import csv
        # openpyxl 시도
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(filepath), read_only=True, data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"## Sheet: {sheet_name}\n")
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    rows.append(" | ".join(cells))
                if rows:
                    header = rows[0]
                    separator = " | ".join(["---"] * len(rows[0].split(" | ")))
                    text_parts.append(f"{header}\n{separator}\n" + "\n".join(rows[1:]))
            wb.close()
            meta = {"format": "XLSX", "sheets": len(wb.sheetnames)}
            return "\n\n".join(text_parts), meta
        except ImportError:
            return "", {"error": "openpyxl not installed"}
    except Exception as e:
        return "", {"error": str(e)}


def write_markdown(category: str, slug: str, content: str, meta: dict, source_path: Path):
    """추출 결과를 마크다운으로 저장"""
    output_dir = OUTPUT_ROOT / category
    output_dir.mkdir(parents=True, exist_ok=True)
    output_file = output_dir / f"{slug}.md"

    file_size_mb = source_path.stat().st_size / (1024 * 1024)
    now = datetime.now().strftime("%Y-%m-%d %H:%M")

    header = f"""---
source: "{source_path}"
category: {category}
format: {meta.get('format', 'unknown')}
file_size: {file_size_mb:.1f}MB
extracted: {now}
"""
    if "total_pages" in meta:
        header += f"total_pages: {meta['total_pages']}\nextracted_pages: {meta['extracted_pages']}\n"
    if "total_slides" in meta:
        header += f"total_slides: {meta['total_slides']}\nextracted_slides: {meta['extracted_slides']}\n"
    if "error" in meta:
        header += f"error: \"{meta['error']}\"\n"
    header += "---\n\n"

    title = f"# {source_path.stem}\n\n"

    with open(output_file, "w", encoding="utf-8") as f:
        f.write(header + title + content)

    return output_file


def is_excluded(filepath: Path) -> bool:
    """개인 문서 제외 체크"""
    name = filepath.name
    return any(pattern in name for pattern in EXCLUDE_PATTERNS)


def process_document(source_rel: str, category: str, slug: str,
                     max_pages: int = 0, max_slides: int = 0,
                     dry_run: bool = False) -> dict:
    """단일 문서 처리"""
    source_path = KRAFTON_ROOT / source_rel

    if not source_path.exists():
        return {"slug": slug, "status": "not_found", "path": str(source_path)}

    if is_excluded(source_path):
        return {"slug": slug, "status": "excluded", "path": str(source_path)}

    file_size_mb = source_path.stat().st_size / (1024 * 1024)
    suffix = source_path.suffix.lower()

    if dry_run:
        return {
            "slug": slug,
            "status": "dry_run",
            "path": str(source_path),
            "size_mb": f"{file_size_mb:.1f}",
            "format": suffix,
            "category": category,
        }

    print(f"  [{suffix.upper()[1:]}] {source_path.name} ({file_size_mb:.1f}MB) → {category}/{slug}.md")

    if suffix == ".pdf":
        content, meta = extract_pdf(source_path, max_pages)
    elif suffix == ".pptx":
        content, meta = extract_pptx(source_path, max_slides)
    elif suffix == ".xlsx":
        content, meta = extract_xlsx(source_path)
    else:
        return {"slug": slug, "status": "unsupported_format", "format": suffix}

    if "error" in meta:
        print(f"    [WARN] Error: {meta['error']}")
        # 에러여도 파일은 생성 (에러 메시지 포함)
        content = f"**추출 실패**: {meta['error']}\n\n수동 처리가 필요합니다."

    output_file = write_markdown(category, slug, content, meta, source_path)
    content_len = len(content)

    return {
        "slug": slug,
        "status": "success" if "error" not in meta else "error",
        "output": str(output_file),
        "size_mb": f"{file_size_mb:.1f}",
        "content_length": content_len,
        "category": category,
    }


def generate_index(results: list[dict]):
    """추출 인덱스 파일 생성"""
    index_path = OUTPUT_ROOT / "_index.md"

    categories = {}
    for r in results:
        if r["status"] in ("success", "error"):
            cat = r.get("category", "unknown")
            if cat not in categories:
                categories[cat] = []
            categories[cat].append(r)

    lines = [
        "# Krafton 문서 추출 인덱스\n",
        f"- **추출 일시**: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        f"- **총 추출 문서**: {sum(len(v) for v in categories.values())}개\n",
    ]

    for cat, items in sorted(categories.items()):
        lines.append(f"\n## {cat} ({len(items)}개)\n")
        for item in items:
            status_icon = "[OK]" if item["status"] == "success" else "[WARN]"
            lines.append(f"- {status_icon} [{item['slug']}]({cat}/{item['slug']}.md) — {item.get('size_mb', '?')}MB")

    with open(index_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    print(f"\n[INDEX] 인덱스 생성: {index_path}")


def main():
    parser = argparse.ArgumentParser(description="D:\\Krafton\\ 문서 텍스트 추출")
    parser.add_argument("--test", action="store_true", help="소형 PDF 1개만 테스트")
    parser.add_argument("--dry-run", action="store_true", help="실제 추출 없이 대상 목록만 표시")
    parser.add_argument("--batch", choices=["small", "medium", "large", "all"], default="all",
                        help="배치 크기: small(<10MB), medium(10-100MB), large(100MB+), all")
    parser.add_argument("--max-pages", type=int, default=0, help="PDF 최대 페이지 수 (0=전체)")
    parser.add_argument("--max-slides", type=int, default=0, help="PPTX 최대 슬라이드 수 (0=전체)")
    parser.add_argument("--single", type=str, help="특정 slug만 추출")
    args = parser.parse_args()

    print("=" * 60)
    print("D:\\Krafton\\ 문서 텍스트 추출기")
    print("=" * 60)

    if args.test:
        # 가장 작은 PDF 1개로 테스트
        test_doc = DOCUMENT_MAP[5]  # 재미와 디자인과 개발 (696K)
        print(f"\n[TEST] 테스트 모드: {test_doc[0]}")
        result = process_document(test_doc[0], test_doc[1], test_doc[2])
        print(f"  결과: {result['status']}")
        if result["status"] == "success":
            print(f"  출력: {result['output']}")
            print(f"  내용 길이: {result['content_length']} chars")
        return

    # 배치 필터링
    results = []
    for source_rel, category, slug in DOCUMENT_MAP:
        if args.single and args.single != slug:
            continue

        source_path = KRAFTON_ROOT / source_rel
        if not source_path.exists():
            results.append({"slug": slug, "status": "not_found", "path": str(source_path)})
            continue

        file_size_mb = source_path.stat().st_size / (1024 * 1024)

        if args.batch == "small" and file_size_mb >= 10:
            continue
        elif args.batch == "medium" and (file_size_mb < 10 or file_size_mb >= 100):
            continue
        elif args.batch == "large" and file_size_mb < 100:
            continue

        result = process_document(
            source_rel, category, slug,
            max_pages=args.max_pages,
            max_slides=args.max_slides,
            dry_run=args.dry_run,
        )
        results.append(result)

    # 결과 요약
    success = sum(1 for r in results if r["status"] == "success")
    errors = sum(1 for r in results if r["status"] == "error")
    not_found = sum(1 for r in results if r["status"] == "not_found")
    excluded = sum(1 for r in results if r["status"] == "excluded")

    print(f"\n{'=' * 60}")
    print(f"결과: [OK] {success} 성공, [WARN] {errors} 에러, [MISS] {not_found} 미발견, [SKIP] {excluded} 제외")

    if not args.dry_run and success > 0:
        generate_index(results)

    # 실패한 항목 표시
    for r in results:
        if r["status"] == "not_found":
            print(f"  [MISS] {r['slug']}: {r['path']}")
        elif r["status"] == "error":
            print(f"  [WARN] {r['slug']}: 추출 에러 (마크다운은 생성됨)")


if __name__ == "__main__":
    main()
